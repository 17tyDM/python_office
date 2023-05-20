from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "ハローワールド!"
subtitle.text = "こんにちは、あなたの健康を守ります。"

prs.save('test2.pptx')